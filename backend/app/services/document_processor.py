"""Document processor — extract text from PDF/TXT/MD files."""

import asyncio
import io
import logging
import os
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


def _extract_text_from_pdf_sync(file_path: Path) -> Optional[list[dict]]:
    """Sync helper: extract text and page boundaries from PDF using PyMuPDF.

    For image-only (scanned) pages where PyMuPDF returns no text,
    falls back to OCR via pytesseract on embedded images.
    """
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        pages = []
        for i in range(len(doc)):
            page = doc[i]
            text = page.get_text("text", flags=0).strip()
            # Fallback to OCR on image-only pages (scanned documents)
            if not text:
                images = page.get_images(full=True)
                if images:
                    try:
                        import tempfile
                        from app.services.ocr_engine import get_ocr_engine

                        xref = images[0][0]
                        base_image = doc.extract_image(xref)
                        img_bytes = base_image["image"]
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                            tmp.write(img_bytes)
                            tmp_path = Path(tmp.name)
                        try:
                            engine = get_ocr_engine()
                            ocr_results = engine.recognize(tmp_path)
                            text = " ".join(r["text"] for r in ocr_results)
                        finally:
                            tmp_path.unlink(missing_ok=True)
                    except Exception as ocr_err:
                        logger.warning(f"OCR failed for page {i + 1}: {ocr_err}")
            pages.append({
                "page_num": i + 1,
                "text": text.strip() if isinstance(text, str) else "",
                "width": page.rect.width,
                "height": page.rect.height,
            })
        doc.close()
        return pages
    except ImportError:
        logger.warning("PyMuPDF not installed — cannot process PDFs")
        return None
    except Exception as e:
        logger.error(f"PDF extraction failed: {e}")
        return None


async def extract_text_from_pdf(file_path: Path) -> Optional[list[dict]]:
    """Extract text and page boundaries from PDF using PyMuPDF.

    For image-only (scanned) pages where PyMuPDF returns no text,
    falls back to OCR via pytesseract on embedded images.
    """
    return await asyncio.to_thread(_extract_text_from_pdf_sync, file_path)


def _extract_text_from_spreadsheet_sync(file_path: Path) -> Optional[str]:
    """Sync helper: extract text from Excel/CSV files using pandas."""
    try:
        import pandas as pd
        if file_path.suffix.lower() == ".csv":
            try:
                df = pd.read_csv(file_path)
            except UnicodeDecodeError:
                df = pd.read_csv(file_path, encoding="latin-1")
        else:
            df = pd.read_excel(file_path)

        # Convert the dataframe to a string representation
        return df.to_string(index=False)
    except ImportError:
        logger.warning("pandas not installed — cannot process spreadsheets")
        return None
    except Exception as e:
        logger.error(f"Spreadsheet extraction failed: {e}")
        return None


async def extract_text_from_spreadsheet(file_path: Path) -> Optional[str]:
    """Extract text from Excel/CSV files using pandas."""
    return await asyncio.to_thread(_extract_text_from_spreadsheet_sync, file_path)


def _extract_text_from_docx_sync(file_path: Path) -> Optional[str]:
    """Sync helper: extract text from Word documents using python-docx."""
    try:
        from docx import Document
        doc = Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])
    except ImportError:
        logger.warning("python-docx not installed — cannot process Word docs")
        return None
    except Exception as e:
        logger.error(f"Docx extraction failed: {e}")
        return None


async def extract_text_from_docx(file_path: Path) -> Optional[str]:
    """Extract text from Word documents using python-docx."""
    return await asyncio.to_thread(_extract_text_from_docx_sync, file_path)


def _extract_text_from_image_sync(file_path: Path) -> Optional[str]:
    """Sync helper: extract text from images using OCR via the ocr_engine abstraction."""
    try:
        from app.services.ocr_engine import get_ocr_engine

        engine = get_ocr_engine()
        results = engine.recognize(file_path)
        text = "\n".join(r["text"] for r in results)
        return text or None
    except Exception as e:
        logger.error(f"Image extraction failed: {e}")
        return None


async def extract_text_from_image(file_path: Path) -> Optional[str]:
    """Extract text from images using OCR via the ocr_engine abstraction."""
    return await asyncio.to_thread(_extract_text_from_image_sync, file_path)


def _extract_text_from_pptx_sync(file_path: Path) -> Optional[str]:
    """Sync helper: extract text from PowerPoint presentations."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except ImportError:
        logger.warning("python-pptx not installed")
        return None
    except Exception as e:
        logger.error(f"PPTX extraction failed: {e}")
        return None


async def extract_text_from_pptx(file_path: Path) -> Optional[str]:
    """Extract text from PowerPoint presentations."""
    return await asyncio.to_thread(_extract_text_from_pptx_sync, file_path)


def _extract_text_from_html_sync(file_path: Path) -> Optional[str]:
    """Sync helper: extract text from HTML files."""
    try:
        from bs4 import BeautifulSoup
        html = file_path.read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(html, "lxml")
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        return soup.get_text(separator="\n")
    except ImportError:
        logger.warning("beautifulsoup4 or lxml not installed")
        return None
    except Exception as e:
        logger.error(f"HTML extraction failed: {e}")
        return None


async def extract_text_from_html(file_path: Path) -> Optional[str]:
    """Extract text from HTML files."""
    return await asyncio.to_thread(_extract_text_from_html_sync, file_path)


def _extract_text_from_odt_sync(file_path: Path) -> Optional[str]:
    """Sync helper: extract text from OpenDocument Text (.odt)."""
    try:
        from odf import text, teletype
        from odf.opendocument import load
        textdoc = load(str(file_path))
        allparagraphs = textdoc.getElementsByType(text.P)
        return "\n".join([teletype.extractText(p) for p in allparagraphs])
    except ImportError:
        logger.warning("odfpy not installed")
        return None
    except Exception as e:
        logger.error(f"ODT extraction failed: {e}")
        return None


async def extract_text_from_odt(file_path: Path) -> Optional[str]:
    """Extract text from OpenDocument Text (.odt)."""
    return await asyncio.to_thread(_extract_text_from_odt_sync, file_path)


def _extract_text_from_rtf_sync(file_path: Path) -> Optional[str]:
    """Sync helper: extract text from RTF files."""
    try:
        from striprtf.striprtf import rtf_to_text
        rtf = file_path.read_text(encoding="utf-8", errors="ignore")
        return rtf_to_text(rtf)
    except ImportError:
        logger.warning("striprtf not installed")
        return None
    except Exception as e:
        logger.error(f"RTF extraction failed: {e}")
        return None


async def extract_text_from_rtf(file_path: Path) -> Optional[str]:
    """Extract text from RTF files."""
    return await asyncio.to_thread(_extract_text_from_rtf_sync, file_path)


def _extract_text_from_epub_sync(file_path: Path) -> Optional[str]:
    """Sync helper: extract text from EPUB files."""
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
        import html2text

        book = epub.read_epub(str(file_path))
        chapters = []
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                chapters.append(item.get_content())

        h = html2text.HTML2Text()
        h.ignore_links = True
        return "\n".join([h.handle(c.decode("utf-8")) for c in chapters])
    except ImportError:
        logger.warning("ebooklib or html2text not installed")
        return None
    except Exception as e:
        logger.error(f"EPUB extraction failed: {e}")
        return None


async def extract_text_from_epub(file_path: Path) -> Optional[str]:
    """Extract text from EPUB files."""
    return await asyncio.to_thread(_extract_text_from_epub_sync, file_path)


def _extract_text_from_email_sync(file_path: Path) -> Optional[str]:
    """Sync helper: extract text from .msg or .eml files."""
    ext = file_path.suffix.lower()
    try:
        if ext == ".msg":
            import extract_msg
            msg = extract_msg.Message(str(file_path))
            return f"Subject: {msg.subject}\nFrom: {msg.sender}\nDate: {msg.date}\n\n{msg.body}"
        elif ext == ".eml":
            import email
            from email.policy import default as default_policy
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                mail = email.message_from_file(f, policy=default_policy)
            subject = mail.get("Subject", "")
            from_addr = mail.get("From", "")
            date = mail.get("Date", "")
            # Extract body text
            body = ""
            if mail.is_multipart():
                for part in mail.walk():
                    if part.get_content_type() == "text/plain":
                        body += part.get_content()
            else:
                body = mail.get_content()
            return f"Subject: {subject}\nFrom: {from_addr}\nDate: {date}\n\n{body}"
    except ImportError:
        logger.warning("extract-msg not installed for .msg support")
        return None
    except Exception as e:
        logger.error(f"Email extraction failed: {e}")
        return None


async def extract_text_from_email(file_path: Path) -> Optional[str]:
    """Extract text from .msg or .eml files."""
    return await asyncio.to_thread(_extract_text_from_email_sync, file_path)


def _extract_zip_to_temp(file_path: Path) -> Optional[Path]:
    """Sync helper: extract a zip to a temp directory (zip-slip safe).

    Returns the temp directory path on success, or None on zip-slip /
    bad-zip / OS error. The temp directory is cleaned up internally on
    any failure path; on success the caller is responsible for cleanup.
    """
    import zipfile
    import tempfile
    import shutil

    tmpdir = None
    success = False
    try:
        tmpdir = Path(tempfile.mkdtemp())
        tmpdir_resolved = tmpdir.resolve()
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            # Zip-slip protection: validate every member path before extraction
            for member in zip_ref.namelist():
                member_path = (tmpdir_resolved / member).resolve()
                if not str(member_path).startswith(str(tmpdir_resolved)):
                    logger.error(f"Zip-slip detected: {member} escapes extraction directory")
                    return None
            zip_ref.extractall(tmpdir)
        success = True
        return tmpdir
    except (zipfile.BadZipFile, OSError) as e:
        logger.error(f"Archive extraction failed: {e}")
        return None
    finally:
        if tmpdir is not None and not success:
            shutil.rmtree(tmpdir, ignore_errors=True)


def _walk_files(root: Path) -> list[Path]:
    """Sync helper: collect all file paths under root."""
    return [Path(r) / f for r, _, fs in os.walk(root) for f in fs]


async def extract_text_from_archive(file_path: Path) -> Optional[str]:
    """Extract text from all files in a ZIP archive with zip-slip protection."""
    import shutil

    tmpdir = await asyncio.to_thread(_extract_zip_to_temp, file_path)
    if tmpdir is None:
        return None

    try:
        files = await asyncio.to_thread(_walk_files, tmpdir)
        text_outputs = []
        for full_p in files:
            # Recursively call extract_text for each file
            res = await extract_text(full_p)
            if res and res.get("content"):
                rel_p = full_p.relative_to(tmpdir)
                text_outputs.append(f"--- FILE: {rel_p} ---\n{res['content']}")
        return "\n\n".join(text_outputs)
    finally:
        await asyncio.to_thread(shutil.rmtree, tmpdir, ignore_errors=True)

async def extract_text(file_path: Path) -> Optional[dict]:
    """Extract text from any supported format. Returns {type, content, pages}."""
    ext = file_path.suffix.lower()
    
    pages = []
    full_text = ""
    
    if ext == ".pdf":
        pages = await extract_text_from_pdf(file_path)
        if pages is None:
            return None
        full_text = "\n".join(p["text"] for p in pages)
        
    elif ext in (".xls", ".xlsx", ".csv"):
        full_text = await extract_text_from_spreadsheet(file_path)
        
    elif ext == ".docx":
        full_text = await extract_text_from_docx(file_path)
            
    elif ext in (".jpg", ".png", ".bmp", ".gif", ".jpeg"):
        full_text = await extract_text_from_image(file_path)

    elif ext == ".pptx":
        full_text = await extract_text_from_pptx(file_path)

    elif ext in (".html", ".htm"):
        full_text = await extract_text_from_html(file_path)

    elif ext == ".odt":
        full_text = await extract_text_from_odt(file_path)

    elif ext == ".rtf":
        full_text = await extract_text_from_rtf(file_path)

    elif ext == ".epub":
        full_text = await extract_text_from_epub(file_path)

    elif ext in (".msg", ".eml"):
        full_text = await extract_text_from_email(file_path)

    elif ext == ".zip":
        full_text = await extract_text_from_archive(file_path)
            
    elif ext in (".txt", ".md", ".py", ".js", ".ts", ".cpp", ".c", ".h", ".java", ".json", ".yaml", ".yml", ".xml"):
        try:
            full_text = file_path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            full_text = file_path.read_text(encoding="latin-1")
            
    else:
        # Fallback: try reading as text if it's not a known binary format
        try:
            full_text = file_path.read_text(encoding="utf-8")
        except (UnicodeDecodeError, ValueError, OSError):
            logger.warning(f"Cannot read as text, unsupported file type: {ext}")
            return None
            
    if full_text is None:
        return None
        
    # Standardize output for non-PDFs
    if ext != ".pdf":
        pages = [{"page_num": 1, "text": full_text}]
        
    return {
        "type": ext.lstrip("."),
        "content": full_text,
        "pages": pages,
        "page_count": len(pages),
    }


def recursive_chunk(text: str, chunk_size: int = 512, overlap: int = 64) -> list[dict]:
    """Split text into chunks with overlap."""
    tokens = text.split()
    chunks = []
    i = 0
    chunk_idx = 0
    while i < len(tokens):
        chunk_tokens = tokens[i:i + chunk_size]
        if len(chunk_tokens) < chunk_size // 2 and chunk_idx > 0:
            break
        chunks.append({
            "chunk_id": chunk_idx,
            "text": " ".join(chunk_tokens),
            "token_count": len(chunk_tokens),
            "start_token": i,
            "end_token": i + len(chunk_tokens),
        })
        i += chunk_size - overlap
        chunk_idx += 1
    return chunks
